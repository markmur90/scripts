import os
import shutil
import tempfile
import zipfile
from pptx import Presentation
import boto3
import xml.etree.ElementTree as ET
from mutagen.mp3 import MP3

ppt_path = "/home/markmur88/scripts/utils/narrar_pptx/presentaciones/presentacion.pptx"
output_ppt = "/home/markmur88/scripts/utils/narrar_pptx/presentaciones/presentacion_narrada_con_tiempos.pptx"
audio_folder = "/home/markmur88/scripts/utils/narrar_pptx/audios"

os.makedirs(audio_folder, exist_ok=True)
polly = boto3.client("polly", region_name="us-east-1")
prs = Presentation(ppt_path)
xml_ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

for idx, slide in enumerate(prs.slides, start=1):
    texts = []
    for shape in slide.shapes:
        elem = shape.element
        for t in elem.findall(".//a:t", xml_ns):
            if t.text:
                texts.append(t.text)
    full_text = "\n".join(texts)
    if not full_text:
        continue
    resp = polly.synthesize_speech(
        Text=full_text,
        VoiceId="Andres",
        OutputFormat="mp3",
        Engine="generative"
    )
    audio_path = os.path.join(audio_folder, f"slide{idx}.mp3")
    with open(audio_path, "wb") as f:
        f.write(resp["AudioStream"].read())

tmpdir = tempfile.mkdtemp()
with zipfile.ZipFile(ppt_path, "r") as zin:
    zin.extractall(tmpdir)

media_dir = os.path.join(tmpdir, "ppt", "media")
os.makedirs(media_dir, exist_ok=True)
for filename in os.listdir(audio_folder):
    shutil.copy(os.path.join(audio_folder, filename), media_dir)

ns = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
}

for idx in range(1, len(prs.slides) + 1):
    rels_path = os.path.join(tmpdir, "ppt", "slides", "_rels", f"slide{idx}.xml.rels")
    rels_tree = ET.parse(rels_path)
    rels_root = rels_tree.getroot()
    rid = f"rId{len(rels_root) + 1}"
    ET.SubElement(
        rels_root,
        "Relationship",
        Id=rid,
        Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/media",
        Target=f"../media/slide{idx}.mp3"
    )
    rels_tree.write(rels_path, xml_declaration=True, encoding="UTF-8", method="xml")

    audio_file = os.path.join(media_dir, f"slide{idx}.mp3")
    duration_ms = int(MP3(audio_file).info.length * 1000)

    slide_xml = os.path.join(tmpdir, "ppt", "slides", f"slide{idx}.xml")
    tree = ET.parse(slide_xml)
    root = tree.getroot()

    spTree = root.find(".//p:spTree", ns)
    if spTree is None:
        continue

    audio = ET.SubElement(spTree, f"{{{ns['p']}}}audio")
    ET.SubElement(audio, f"{{{ns['p']}}}cTn", {"id": "0", "dur": "0", "fill": "hold"})
    ET.SubElement(audio, f"{{{ns['p']}}}audioFile", {f"{{{ns['r']}}}link": rid})

    for old in root.findall("p:transition", ns):
        root.remove(old)
    ET.SubElement(
        root,
        f"{{{ns['p']}}}transition",
        {"advTm": str(duration_ms), "advClick": "0"}
    )

    tree.write(slide_xml, xml_declaration=True, encoding="UTF-8", method="xml")

ct_path = os.path.join(tmpdir, "[Content_Types].xml")
ct_tree = ET.parse(ct_path)
ct_root = ct_tree.getroot()
for idx in range(1, len(prs.slides) + 1):
    ET.SubElement(
        ct_root,
        "Override",
        PartName=f"/ppt/media/slide{idx}.mp3",
        ContentType="audio/mpeg"
    )
ct_tree.write(ct_path, xml_declaration=True, encoding="UTF-8", method="xml")

with zipfile.ZipFile(output_ppt, "w", zipfile.ZIP_DEFLATED) as zout:
    for dirpath, _, files in os.walk(tmpdir):
        for f in files:
            full = os.path.join(dirpath, f)
            arc = os.path.relpath(full, tmpdir)
            zout.write(full, arc)

shutil.rmtree(tmpdir)
