import os
import shutil
import tempfile
import subprocess
from glob import glob
from pptx import Presentation
import boto3
from mutagen.mp3 import MP3

base = os.path.expanduser("/home/markmur88/scripts/utils/narrar_pptx")
presentaciones_dir = os.path.join(base, "presentaciones")
audio_folder = os.path.join(base, "audios")
img_dir = os.path.join(base, "img")
videos_dir = os.path.join(base, "videos")
os.makedirs(presentaciones_dir, exist_ok=True)
os.makedirs(audio_folder, exist_ok=True)
os.makedirs(img_dir, exist_ok=True)
os.makedirs(videos_dir, exist_ok=True)
ppt_filename = "presentacion.pptx"
ppt_path = os.path.join(presentaciones_dir, ppt_filename)
if not os.path.isfile(ppt_path):
    raise FileNotFoundError(f"No se encuentra PPTX en: {ppt_path}")
libre = shutil.which("libreoffice") or shutil.which("soffice")
if not libre:
    raise RuntimeError("No se encontró libreoffice/soffice en PATH")
subprocess.run([libre, "--headless", "--convert-to", "pdf:writer_pdf_Export", "--outdir", presentaciones_dir, ppt_path], check=True)
pdf_path = os.path.join(presentaciones_dir, os.path.splitext(ppt_filename)[0] + ".pdf")
if not os.path.isfile(pdf_path):
    raise RuntimeError(f"No se generó PDF en: {pdf_path}")
subprocess.run(["convert", "-density", "150", pdf_path, os.path.join(img_dir, "slide_%03d.png")], check=True)
polly = boto3.client("polly", region_name="us-east-1")
prs = Presentation(ppt_path)
for idx, slide in enumerate(prs.slides, start=1):
    textos = []
    for shp in slide.shapes:
        for t in shp.element.findall(".//a:t", {"a":"http://schemas.openxmlformats.org/drawingml/2006/main"}):
            if t.text:
                textos.append(t.text)
    contenido = "\n".join(textos)
    if not contenido:
        continue
    resp = polly.synthesize_speech(Text=contenido, VoiceId="Andres", OutputFormat="mp3", Engine="generative")
    ruta_mp3 = os.path.join(audio_folder, f"slide{idx}.mp3")
    with open(ruta_mp3, "wb") as f:
        f.write(resp["AudioStream"].read())
slides = sorted(glob(os.path.join(img_dir, "slide_*.png")))
audios = sorted(glob(os.path.join(audio_folder, "slide*.mp3")))
n = min(len(slides), len(audios))
if n == 0:
    raise RuntimeError("No hay imágenes o audios para procesar")
tmp = tempfile.mkdtemp()
segmentos = []
for i in range(n):
    img = slides[i]
    aud = audios[i]
    seg = os.path.join(tmp, f"seg{i+1:03d}.mp4")
    subprocess.run([
        "ffmpeg", "-y", "-loop", "1", "-i", img, "-i", aud,
        "-vf", "scale=if(gt(mod(iw,2),0),iw+1,iw):if(gt(mod(ih,2),0),ih+1,ih)",
        "-c:v", "libx264", "-c:a", "aac", "-b:a", "192k",
        "-pix_fmt", "yuv420p", "-shortest", seg
    ], check=True)
    segmentos.append(seg)
concat_txt = os.path.join(tmp, "concat.txt")
with open(concat_txt, "w") as f:
    for seg in segmentos:
        f.write(f"file '{seg}'\n")
final_video = os.path.join(videos_dir, "presentacion_final.mp4")
subprocess.run(["ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", concat_txt, "-c", "copy", final_video], check=True)
shutil.rmtree(tmp)
print(f"Vídeo generado en:\n  {final_video}")
